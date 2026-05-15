import io
import pandas as pd
from datetime import datetime

# PDF
from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    PageBreak,
    Table,
    TableStyle,
)
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib.pagesizes import landscape, A4

# POWERPOINT
from pptx import Presentation
from pptx.util import Inches


# ==========================================================
# HELPERS
# ==========================================================
def safe_df(df):
    if df is None:
        return pd.DataFrame()

    if isinstance(df, pd.DataFrame):
        return df.copy()

    return pd.DataFrame()


def add_dataframe_to_pdf_story(df, story, title):
    styles = getSampleStyleSheet()

    story.append(Paragraph(f"<b>{title}</b>", styles["Heading2"]))
    story.append(Spacer(1, 12))

    if df.empty:
        story.append(Paragraph("No data available.", styles["BodyText"]))
        story.append(Spacer(1, 12))
        return

    max_rows = min(len(df), 20)

    table_data = [list(df.columns)]

    for _, row in df.head(max_rows).iterrows():
        table_data.append([str(x)[:50] for x in row.tolist()])

    table = Table(table_data)

    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0B1F3A")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("GRID", (0, 0), (-1, -1), 1, colors.grey),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 8),
                ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
                ("BACKGROUND", (0, 1), (-1, -1), colors.whitesmoke),
            ]
        )
    )

    story.append(table)
    story.append(Spacer(1, 20))


# ==========================================================
# PDF REPORT GENERATOR
# ==========================================================
def generate_pdf_report(
    company_name,
    executive_summary,
    total_revenue_opportunity,
    total_leakage_exposure,
    total_cost_saving,
    business_diagnostic_df=None,
    revenue_opportunity_df=None,
    leakage_df=None,
    cost_saving_df=None,
    action_plan_df=None,
    recommendations_df=None,
):
    buffer = io.BytesIO()

    doc = SimpleDocTemplate(
        buffer,
        pagesize=landscape(A4),
        rightMargin=20,
        leftMargin=20,
        topMargin=20,
        bottomMargin=20,
    )

    styles = getSampleStyleSheet()
    story = []

    # ======================================================
    # COVER PAGE
    # ======================================================
    story.append(
        Paragraph(
            f"<font size=24><b>{company_name}</b></font>",
            styles["Title"],
        )
    )

    story.append(Spacer(1, 20))

    story.append(
        Paragraph(
            "<font size=18><b>ProfitIQ Executive Business Review Report</b></font>",
            styles["Title"],
        )
    )

    story.append(Spacer(1, 20))

    story.append(
        Paragraph(
            f"Generated on: {datetime.now().strftime('%d %B %Y')}",
            styles["BodyText"],
        )
    )

    story.append(PageBreak())

    # ======================================================
    # EXECUTIVE SUMMARY
    # ======================================================
    story.append(
        Paragraph(
            "<b>Executive Summary</b>",
            styles["Heading1"],
        )
    )

    story.append(Spacer(1, 12))

    story.append(
        Paragraph(executive_summary, styles["BodyText"])
    )

    story.append(Spacer(1, 20))

    # KPI SUMMARY TABLE
    kpi_df = pd.DataFrame(
        {
            "Metric": [
                "Revenue Opportunity",
                "Leakage Exposure",
                "Cost Saving Potential",
                "Combined Improvement Potential",
            ],
            "Value": [
                f"₦{total_revenue_opportunity:,.0f}",
                f"₦{total_leakage_exposure:,.0f}",
                f"₦{total_cost_saving:,.0f}",
                f"₦{(total_revenue_opportunity + total_leakage_exposure + total_cost_saving):,.0f}",
            ],
        }
    )

    add_dataframe_to_pdf_story(kpi_df, story, "Executive KPI Summary")

    # ======================================================
    # REPORT SECTIONS
    # ======================================================
    add_dataframe_to_pdf_story(
        safe_df(business_diagnostic_df),
        story,
        "Business Diagnostic Report",
    )

    add_dataframe_to_pdf_story(
        safe_df(revenue_opportunity_df),
        story,
        "Revenue Opportunity Map",
    )

    add_dataframe_to_pdf_story(
        safe_df(leakage_df),
        story,
        "Leakage Register",
    )

    add_dataframe_to_pdf_story(
        safe_df(cost_saving_df),
        story,
        "Cost Savings Report",
    )

    add_dataframe_to_pdf_story(
        safe_df(action_plan_df),
        story,
        "90-Day Profit Improvement Plan",
    )

    add_dataframe_to_pdf_story(
        safe_df(recommendations_df),
        story,
        "Management Recommendations",
    )

    doc.build(story)

    pdf = buffer.getvalue()
    buffer.close()

    return pdf


# ==========================================================
# POWERPOINT GENERATOR
# ==========================================================
def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_text_slide(prs, title, content):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = content


def add_table_slide(prs, title, df):
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title

    df = safe_df(df)

    if df.empty:
        textbox = slide.shapes.add_textbox(
            Inches(1),
            Inches(2),
            Inches(8),
            Inches(1),
        )

        textbox.text = "No data available."
        return

    rows = min(len(df), 10) + 1
    cols = len(df.columns)

    table = slide.shapes.add_table(
        rows,
        cols,
        Inches(0.5),
        Inches(1.5),
        Inches(12),
        Inches(4),
    ).table

    for col_idx, col_name in enumerate(df.columns):
        table.cell(0, col_idx).text = str(col_name)

    for row_idx, (_, row) in enumerate(df.head(10).iterrows(), start=1):
        for col_idx, value in enumerate(row):
            table.cell(row_idx, col_idx).text = str(value)[:40]


def generate_powerpoint_report(
    company_name,
    executive_summary,
    total_revenue_opportunity,
    total_leakage_exposure,
    total_cost_saving,
    revenue_opportunity_df=None,
    leakage_df=None,
    cost_saving_df=None,
    action_plan_df=None,
):
    prs = Presentation()

    # ======================================================
    # TITLE SLIDE
    # ======================================================
    add_title_slide(
        prs,
        f"{company_name}",
        "ProfitIQ Executive Business Review",
    )

    # ======================================================
    # EXECUTIVE SUMMARY
    # ======================================================
    add_text_slide(
        prs,
        "Executive Summary",
        executive_summary,
    )

    # ======================================================
    # KPI SUMMARY
    # ======================================================
    kpi_content = f"""
Revenue Opportunity: ₦{total_revenue_opportunity:,.0f}

Leakage Exposure: ₦{total_leakage_exposure:,.0f}

Cost Saving Potential: ₦{total_cost_saving:,.0f}

Combined Profit Improvement Potential:
₦{(total_revenue_opportunity + total_leakage_exposure + total_cost_saving):,.0f}
"""

    add_text_slide(
        prs,
        "Executive KPI Summary",
        kpi_content,
    )

    # ======================================================
    # TABLE SLIDES
    # ======================================================
    add_table_slide(
        prs,
        "Revenue Opportunity Map",
        revenue_opportunity_df,
    )

    add_table_slide(
        prs,
        "Leakage Register",
        leakage_df,
    )

    add_table_slide(
        prs,
        "Cost Savings Report",
        cost_saving_df,
    )

    add_table_slide(
        prs,
        "90-Day Profit Improvement Plan",
        action_plan_df,
    )

    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)

    ppt_data = ppt_buffer.getvalue()
    ppt_buffer.close()

    return ppt_data